from pptx import Presentation
from pptx.util import Inches


prs = Presentation("/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Double-Vorlage.pptx")
karte = ['116.png', '104.png', '123.png', '31.png', '19.png', '56.png', '43.png', '80.png', '55.png', '7.png',
         '68.png', '92.png']
for ind_slide, slide in enumerate(prs.slides):

    top = 0.4
    left = Inches(-0.55)
    height = Inches(1.7)
    try:
        img_path = f"/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Symbolordner/{karte[ind_slide]}"
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
    except:
        pass
prs.save('/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/test.pptx')
##for symbol in karte:
    #symbol_path = f"/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Symbolordner/{symbol}"
