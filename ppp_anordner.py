import os
from pptx import Presentation
from pptx.util import Inches
prs = Presentation("/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Double-Vorlage.pptx")
karten = os.listdir("/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Kartenordner")

top = 0.4
left = Inches(-0.55)
height = Inches(1.7)

for ind_slide, slide in enumerate(prs.slides):
    if ind_slide == len(karten):
        break
    if "karte" in karten[ind_slide].lower():
        card_content = os.listdir(
            f"/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/Kartenordner/{karten[ind_slide]}")
        for symbols in card_content:
            symbol_path = f"/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble" \
                          f"/Kartenordner/{karten[ind_slide]}/{symbols}"
            pic = slide.shapes.add_picture(symbol_path, left, top, height=height)

prs.save('/Users/marvin/Documents/Sonstiges/Spiele/TribbleTroubleDouble/test.pptx')


